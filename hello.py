from flask import Flask, jsonify, request
from flask_cors import CORS
from pptx import Presentation	

app = Flask(__name__)
CORS(app)


@app.route('/api/data')
def api_data():
    data = {
        'message': 'Hello from Flask REST API!',
        'data': 'Some data to be sent to the Angular frontend'
    }
    return jsonify(data)

@app.route('/api/pptx', methods=['POST'])
def api_pptx():
    messages = request.get_json()
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    
    """ Ref for slide types: 
0 ->  title and subtitle
1 ->  title and content
2 ->  section header
3 ->  two content
4 ->  Comparison
5 ->  Title only 
6 ->  Blank
7 ->  Content with caption
8 ->  Pic with caption
"""
  
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save('test.pptx')

    data = {
        'message': 'Pptx is ready',
        'data': 'test.pptx'
    }
    return jsonify(messages)

if __name__ == '__main__':
    app.run()
